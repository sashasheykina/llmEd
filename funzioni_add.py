from datetime import datetime
import json
import streamlit as st
from pptx import Presentation
from io import BytesIO

def fasePreattiva(its, grado, indirizzo, linee_guida, disciplina, uda, periodo, ore_anno, ore_sett, lezione,  contesto, client, docente):

    elencoUda = open("elencoUDA.json", "r").read()
    dettaglioUda = open("dettaglioUDA.json", "r").read()
    formato_output = open("", "r").read()
    elencoLezioni = open("elencoLezioni.json", "r").read()
    current_dateTime = datetime.now()
    date = str(current_dateTime.day) + "/" + str(current_dateTime.month) + "/" + str(current_dateTime.year)
    metodo = "Nella progettazione di una lezione, considera di includere le metodologie didattiche per l'insegnamento e utilizzo di sussidi didattici dove opportuno. "
    role_system = "Agisci da esperto docente di Informatica negli " + its + ". \ " + linee_guida + ". Il tuo compito è fornire la fase preattiva della lezione che ti fornirò nell’ambito della progettazione didattica per la disciplina  " + disciplina + ", in linea con le indicazioni nazionali. \
        Dovrai fornire la progettazione tramite un JSON strutturato come segue: " + formato_output
    contesto_classe = "### ANALISI PROFILO CLASSE \ " + contesto + "### ANNO DI RIFERIMENTO \ " + grado + " ad indirizzo " + indirizzo + ". \ PROGETTAZIONE: \ "
    role_user = "Io ti fornirò l'analisi del profilo della classe e l'anno di riferimento, l'elenco delle UDA, la UDA di riferimento, deattaglio UDA, periodo, \
                monte ore settimanele, monte ore annuale, elenco lezioni, lezione di riferimento, il metodo, nome docente, data di oggi e tuo restituirai la fase preattiva della lezione per la Progettazione Disciplinare nel formato richiesto. \ " \
                + contesto_classe + "### ELENCO UNITA’ DIDATTICA \  " + elencoUda + "### UNITA’ DIDATTICA \ " + uda + "### DETTAGLIO UNITA’ DIDATTICA \ " + dettaglioUda + " ### PERIODO \ " + periodo + " ### NUMERO DI ORE SETTIMANALI \ " + str(ore_sett) + \
                "### TOTALE ORE ANNO SCOLASTICO \ " + str(ore_anno) + \
                "ELENCO LEZIONI: \ " + elencoLezioni + "### LEZIONE \ " + lezione + "### METODO \ " + metodo + "### NOME DOCENTE \ " + docente + "### DATA \ " + date + \
                "LEZIONE FASE PREATTIVA: \ "

    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": role_system},
            {"role": "user", "content": role_user}
        ]
    )

    # Estrarre la risposta testuale dall'oggetto completion
    response_text = completion.choices[0].message.content

    # Pulire la risposta per trovare solo il JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1

    if json_start != -1 and json_end != -1:
        json_text = response_text[json_start:json_end]

        try:
            response_data = json.loads(json_text)
            # Salvataggio della risposta in un file JSON
            elenco_path = ''
            with open(elenco_path, 'w', encoding='utf-8') as json_file:
                json.dump(response_data, json_file, ensure_ascii=False, indent=2)

            print(f"La risposta è stata salvata in {elenco_path}")
        except json.JSONDecodeError as e:
            print("Errore nel decodificare il JSON:", e)
            print("Contenuto JSON estratto:", json_text)
    else:
        print("Errore: non è stato possibile individuare un blocco JSON valido nella risposta.")
        print("Contenuto della risposta:", response_text)


def dettaglioFasePreattiva(its, grado, indirizzo, linee_guida, disciplina, uda, periodo, ore_anno, ore_sett, lezione,  contesto, client, sussidi_didattici):

    elencoUda = open("elencoUDA.json", "r").read()
    dettaglioUda = open("dettaglioUDA.json", "r").read()
    formato_output = open("", "r").read()
    elencoLezioni = open("elencoLezioni.json", "r").read()
    fasePreattivaLezione = open("", "r").read()
    metodo = "Nella progettazione di una lezione, considera di includere strumenti specifici come " + sussidi_didattici + " dove opportuno. "
    durata = "2 ore"
    role_system = "Agisci da esperto docente di Informatica negli " + its + ". \ " + linee_guida + ". Il tuo compito è fornire il dettaglio della fase preattiva della lezione che ti fornirò nell’ambito della progettazione didattica per la disciplina " + disciplina + ", in linea con le indicazioni nazionali. \
        Dovrai fornire la progettazione tramite un JSON strutturato come segue: " + formato_output
    contesto_classe = "### ANALISI PROFILO CLASSE \ " + contesto + "### ANNO DI RIFERIMENTO \ " + grado + " ad indirizzo " + indirizzo + ". \ PROGETTAZIONE: \ "
    role_user = "Io ti fornirò l'analisi del profilo della classe e l'anno di riferimento, l'elenco delle UDA, la UDA di riferimento, deattaglio UDA, periodo, monte ore settimanele, monte ore annuale, elenco lezioni, lezione, fase preattiva della lezione, durata di tutte le fasi della lezione, metodo \
                e tuo restituirai il dettaglio della fase preattiva della lezione per la Progettazione Disciplinare nel formato richiesto. \ " \
                + contesto_classe + "### ELENCO UNITA’ DIDATTICA \  " + elencoUda + "### UNITA’ DIDATTICA \ " + uda + "### DETTAGLIO UNITA’ DIDATTICA \ " + dettaglioUda + " ### PERIODO \ " + periodo + " ### NUMERO DI ORE SETTIMANALI \ " + str(ore_sett) + \
                "### TOTALE ORE ANNO SCOLASTICO \ " + str(ore_anno) + \
                "ELENCO LEZIONI: \ " + elencoLezioni + "### LEZIONE \ " + lezione + "### FASE PREATTIVA LEZIONE \ " + fasePreattivaLezione + "### METODO \ " + metodo + "### DURATA DI TUTTE LE FASI DELLA LIZIONE" + durata +\
                "DETTAGLIO LEZIONE FASE PREATTIVA: \ "

    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": role_system},
            {"role": "user", "content": role_user}
        ]
    )

    # Estrarre la risposta testuale dall'oggetto completion
    response_text = completion.choices[0].message.content

    # Pulire la risposta per trovare solo il JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1

    if json_start != -1 and json_end != -1:
        json_text = response_text[json_start:json_end]

        try:
            response_data = json.loads(json_text)
            # Salvataggio della risposta in un file JSON
            elenco_path = ''
            with open(elenco_path, 'w', encoding='utf-8') as json_file:
                json.dump(response_data, json_file, ensure_ascii=False, indent=2)

            print(f"La risposta è stata salvata in {elenco_path}")
        except json.JSONDecodeError as e:
            print("Errore nel decodificare il JSON:", e)
            print("Contenuto JSON estratto:", json_text)
    else:
        print("Errore: non è stato possibile individuare un blocco JSON valido nella risposta.")
        print("Contenuto della risposta:", response_text)


def faseAttiva(its, grado, indirizzo, linee_guida, disciplina, uda, periodo, ore_anno, ore_sett, lezione,  contesto, client, docente):

    elencoUda = open("elencoUDA.json", "r").read()
    dettaglioUda = open("dettaglioUDA.json", "r").read()
    formato_output = open("", "r").read()
    elencoLezioni = open("elencoLezioni.json", "r").read()
    role_system = "Agisci da esperto docente di Informatica negli "+its +". \ "+linee_guida+". Il tuo compito è fornire il dettaglio dell’Unità didattica di Apprendimento che ti fornirò nell’ambito della progettazione didattica per la disciplina " + disciplina +", in linea con le indicazioni nazionali. \
    Dovrai fornire la progettazione tramite un JSON strutturato come segue: " + formato_output
    contesto_classe = "### ANALISI PROFILO CLASSE \ " + contesto + "### ANNO DI RIFERIMENTO \ " + grado + " ad indirizzo " + indirizzo + ". \ PROGETTAZIONE: \ "
    role_user = "Io ti fornirò l'analisi del profilo della classe e l'anno di riferimento, l'elenco delle UDA, la UDA di riferimento, deattaglio UDA, periodo, monte ore settimanele, monte ore annuale, elenco lezioni, lezione, nome docente, \
                e tuo restituirai la fase preattiva della lezione per la Progettazione Disciplinare nel formato richiesto. \ " \
                + contesto_classe + "### ELENCO UNITA’ DIDATTICA \  " + elencoUda + "### UNITA’ DIDATTICA \ " + uda + "### DETTAGLIO UNITA’ DIDATTICA \ " + dettaglioUda + " ### PERIODO \ " + periodo + " ### NUMERO DI ORE SETTIMANALI \ " + ore_sett + \
                "### TOTALE ORE ANNO SCOLASTICO \ " + ore_anno + \
                "ELENCO LEZIONI: \ " + elencoLezioni+" ### LEZIONE \ "+ lezione + "### NOME DOCENTE: \ " + docente + \
                "LEZIONE FASE ATTIVA: \ "

    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": role_system},
            {"role": "user", "content": role_user}
        ]
    )

    # Estrarre la risposta testuale dall'oggetto completion
    response_text = completion.choices[0].message.content

    # Pulire la risposta per trovare solo il JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1

    if json_start != -1 and json_end != -1:
        json_text = response_text[json_start:json_end]

        try:
            response_data = json.loads(json_text)
            # Salvataggio della risposta in un file JSON
            elenco_path = ''
            with open(elenco_path, 'w', encoding='utf-8') as json_file:
                json.dump(response_data, json_file, ensure_ascii=False, indent=2)

            print(f"La risposta è stata salvata in {elenco_path}")
        except json.JSONDecodeError as e:
            print("Errore nel decodificare il JSON:", e)
            print("Contenuto JSON estratto:", json_text)
    else:
        print("Errore: non è stato possibile individuare un blocco JSON valido nella risposta.")
        print("Contenuto della risposta:", response_text)

def genSlidesFasePreattiva(tipo_scuola, grado, address, linee_guida, disciplina, uda, periodo, ore_anno, ore_sett, lezione, contesto, client, docente):

    elencoUda = open("elencoUDA.json", "r").read()
    dettaglioUda = open("dettaglioUDA.json", "r").read()
    formato_output = open("", "r").read()
    elencoLezioni = open("elencoLezioni.json", "r").read()
    fasePreattiva = open("", "r").read()
    dettaglioFasePreattiva = open("", "r").read()
    current_dateTime = datetime.now()
    date = str(current_dateTime.day) + "/" + str(current_dateTime.month) + "/" + str(current_dateTime.year)
    role_system = "Agisci da esperto docente di Informatica negli "+tipo_scuola +". \ "+linee_guida+". Il  tuo compito è fornire le slides della fase preattiva \
                della lezione che ti fornirò nell’ambito della progettazione didattica per la disciplina " + disciplina +", in linea con le indicazioni nazionali. \
                Dovrai fornire la progettazione tramite un JSON strutturato come segue: " + formato_output
    contesto_classe = "### ANALISI PROFILO CLASSE \ " + contesto + "### ANNO DI RIFERIMENTO \ " + grado + " ad indirizzo " + address +". \ PROGETTAZIONE: \ "
    role_user = "Io ti fornirò l'analisi del profilo della classe e l'anno di riferimento, l'elenco delle UDA, la UDA di riferimento, deattaglio UDA, periodo, \
                monte ore settimanele, monte ore annuale, elenco lezioni, lezione di riferimento, fase preattiva della lezione, dettaglio fase preattiva della lezione, nome docente, data di oggi e tuo restituirai le slides della fase attiva della lezione per la Progettazione Disciplinare nel formato richiesto. \ " \
                + contesto_classe + "### ELENCO UNITA’ DIDATTICA \  " + elencoUda + "### UNITA’ DIDATTICA \ " + uda + "### DETTAGLIO UNITA’ DIDATTICA \ " + dettaglioUda + " ### PERIODO \ " + periodo + " ### NUMERO DI ORE SETTIMANALI \ " + str(ore_sett) + \
                "### TOTALE ORE ANNO SCOLASTICO \ " + str(ore_anno) + \
                "ELENCO LEZIONI \ " + elencoLezioni + "### LEZIONE \ " + lezione + "### NOME DOCENTE \ " + docente + "### DATA \ " + date + \
                "### FASE PREATTIVA \ " + fasePreattiva + "### DETTAGLIO FASE PREATTIVA \ " + dettaglioFasePreattiva + "SLIDES \ "

    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": role_system},
            {"role": "user", "content": role_user}
        ]
    )

    response_text = completion.choices[0].message.content

    # Pulire la risposta per trovare solo il JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1

    if json_start != -1 and json_end != -1:
        json_text = response_text[json_start:json_end]

        try:
            response_data = json.loads(json_text)
            # Salvataggio della risposta in un file JSON
            elenco_path = ''
            with open(elenco_path, 'w', encoding='utf-8') as json_file:
                json.dump(response_data, json_file, ensure_ascii=False, indent=2)

            print(f"La risposta è stata salvata in {elenco_path}")
        except json.JSONDecodeError as e:
            print("Errore nel decodificare il JSON:", e)
            print("Contenuto JSON estratto:", json_text)
    else:
        print("Errore: non è stato possibile individuare un blocco JSON valido nella risposta.")
        print("Contenuto della risposta:", response_text)

def genera_PPT_FasePreattiva():
    # Carica il file JSON
    with open('', 'r') as f:
        fase = json.load(f)

    with open('', 'r') as f:
        dettaglio = json.load(f)

        # Crea una nuova presentazione PowerPoint
        presentation = Presentation()

        # Aggiungi titolo della presentazione
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        subtitle1 = slide.placeholders[1]

        title.text = fase['lesson_title']
        subtitle.text = f"{fase['indirizzo']}\n{fase['grado']}\n{fase['instructor']} - {fase['date']}"

        # Itera attraverso le slide nel JSON
        for slide_data in fase['slides']:
            slide_number = slide_data.get('slide_number')
            title = slide_data.get('title')
            subtitle = slide_data.get('subtitle', '')  # Alcune slide potrebbero non avere sottotitolo
            content = slide_data.get('content', [])
            presenter_notes = slide_data.get('presenter_notes', '')

            # Aggiungi una slide alla presentazione
            slide_layout = presentation.slide_layouts[1]  # Layout della slide (1 è per il titolo e contenuto)
            slide = presentation.slides.add_slide(slide_layout)

            # Aggiungi il titolo alla slide
            title_placeholder = slide.shapes.title
            title_placeholder.text = title

            # Aggiungi il sottotitolo (se presente)
            if subtitle:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = subtitle

            # Aggiungi il contenuto alla slide
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame

            for item in content:
                for key, value in item.items():
                    p = text_frame.add_paragraph()
                    p.text = value

            # Aggiungi le note del presentatore (se presenti)
            if presenter_notes:
                notes_slide = slide.notes_slide
                notes_text_frame = notes_slide.notes_text_frame
                notes_text_frame.text = presenter_notes

        # Salva la presentazione PowerPoint
        presentation.save('output_presentation.pptx')

        print("Presentazione creata con successo e salvata come 'output_presentation.pptx'")

        # save the output into binary form
        ppt_bytes = convert_ppt_to_bytes(presentation)

        # Usa session_state per memorizzare il buffer PPTX
        st.session_state.ppt_bytes = ppt_bytes

# Funzione per convertire il file PPTX in un buffer di byte per il download
def convert_ppt_to_bytes(prs):
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes

def genera_PowerPoint():


    # Carica il file JSON
    with open('slides.json', 'r') as f:
        data = json.load(f)

    # Crea una nuova presentazione PowerPoint
    presentation = Presentation()

    # Aggiungi titolo della presentazione
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    subtitle1 = slide.placeholders[1]

    title.text = data['lesson_title']
    subtitle.text = f"{data['indirizzo']}\n{data['grado']}\n{data['instructor']} - {data['date']}"


    # Itera attraverso le slide nel JSON
    for slide_data in data['slides']:
        slide_number = slide_data.get('slide_number')
        title = slide_data.get('title')
        subtitle = slide_data.get('subtitle', '')  # Alcune slide potrebbero non avere sottotitolo
        content = slide_data.get('content', [])
        presenter_notes = slide_data.get('presenter_notes', '')

        # Aggiungi una slide alla presentazione
        slide_layout = presentation.slide_layouts[1]  # Layout della slide (1 è per il titolo e contenuto)
        slide = presentation.slides.add_slide(slide_layout)

        # Aggiungi il titolo alla slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Aggiungi il sottotitolo (se presente)
        if subtitle:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        # Aggiungi il contenuto alla slide
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame

        for item in content:
            for key, value in item.items():
                p = text_frame.add_paragraph()
                p.text = value

        # Aggiungi le note del presentatore (se presenti)
        if presenter_notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = presenter_notes

    # Salva la presentazione PowerPoint
    presentation.save('output_presentation.pptx')

    print("Presentazione creata con successo e salvata come 'output_presentation.pptx'")

    # save the output into binary form
    ppt_bytes = convert_ppt_to_bytes(presentation)

    # Usa session_state per memorizzare il buffer PPTX
    st.session_state.ppt_bytes = ppt_bytes


def genera_Lezione(tipo_scuola, grado, address, linee_guida, disciplina, uda, periodo, ore_anno, ore_sett, lezione, contesto, client, docente):

    elencoUda = open("elencoUDA.json", "r").read()
    dettaglioUda = open("dettaglioUDA.json", "r").read()
    elencoLezioni = open("elencoLezioni.json", "r").read()
    current_dateTime = datetime.now()
    metodologie = open("metodologie_didattiche", "r").read()
    strumenti = open("materiale_didattico.json", "r").read()
    date = str(current_dateTime.day) + "/" + str(current_dateTime.month) + "/" + str(current_dateTime.year)
    role_system = (
        f"Agisci da esperto docente di Informatica negli {tipo_scuola}. {linee_guida}."
                  f"Il  tuo compito è fornire lezione che ti fornirò nell’ambito della progettazione didattica per la disciplina {disciplina}, in linea con le indicazioni nazionali."
                  f"Dovrai fornire la progettazione tramite un JSON."
            )
    contesto_classe = (
        f"### ANALISI PROFILO CLASSE {contesto} ### ANNO DI RIFERIMENTO {grado} ad indirizzo {address} ### PROGETTAZIONE: "
            )
    role_user = (
        f"Io ti fornirò l'analisi del profilo della classe e l'anno di riferimento, l'elenco delle UDA, la UDA di riferimento, deattaglio UDA, periodo,"
        f"monte ore settimanele, monte ore annuale, elenco lezioni, lezione di riferimento, nome docente, data di oggi, il materiale didattico da utilizzare per la metodologia applicata e tuo restituirai la lezione per la Progettazione Disciplinare nel formato richiesto. {contesto_classe}"
        f"### ELENCO UNITA’ DIDATTICA \  {elencoUda} ### UNITA’ DIDATTICA {uda} ### DETTAGLIO UNITA’ DIDATTICA {dettaglioUda} ### PERIODO {periodo} ### NUMERO DI ORE SETTIMANALI {str(ore_sett)}"
        f"### TOTALE ORE ANNO SCOLASTICO {str(ore_anno)} ###ELENCO LEZIONI {elencoLezioni} ### LEZIONE {lezione} ### NOME DOCENTE {docente} ### DATA {date} ### METODOLOGIE {metodologie} "
        f"### MATERIALE DIDATTICO {strumenti}."
        f"LEZIONE:"
    )

    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": role_system},
            {"role": "user", "content": role_user}
        ]
    )

    response_text = completion.choices[0].message.content

    # Pulire la risposta per trovare solo il JSON
    json_start = response_text.find('{')
    json_end = response_text.rfind('}') + 1

    if json_start != -1 and json_end != -1:
        json_text = response_text[json_start:json_end]

        try:
            response_data = json.loads(json_text)
            # Salvataggio della risposta in un file JSON
            elenco_path = ''
            with open(elenco_path, 'w', encoding='utf-8') as json_file:
                json.dump(response_data, json_file, ensure_ascii=False, indent=2)

            print(f"La risposta è stata salvata in {elenco_path}")
        except json.JSONDecodeError as e:
            print("Errore nel decodificare il JSON:", e)
            print("Contenuto JSON estratto:", json_text)
    else:
        print("Errore: non è stato possibile individuare un blocco JSON valido nella risposta.")
        print("Contenuto della risposta:", response_text)