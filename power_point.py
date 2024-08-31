import json
from pptx import Presentation
from pptx.util import Inches, Pt
import streamlit as st
from io import BytesIO

def genera_PowerPoint():


    # Carica il file JSON
    with open('slides.json', 'r') as f:
        data = json.load(f)

    # Crea una nuova presentazione PowerPoint
    presentation = Presentation()

    # Aggiungi titolo della presentazione
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    subtitle1 = slide.placeholders[1]

    title.text = data['lesson_title']
    subtitle.text = f"{data['indirizzo']}\n{data['grado']}\n{data['instructor']} - {data['date']}"


    # Itera attraverso le slide nel JSON
    for slide_data in data['slides']:
        slide_number = slide_data.get('slide_number')
        title = slide_data.get('title')
        subtitle = slide_data.get('subtitle', '')  # Alcune slide potrebbero non avere sottotitolo
        content = slide_data.get('content', [])
        presenter_notes = slide_data.get('presenter_notes', '')

        # Aggiungi una slide alla presentazione
        slide_layout = presentation.slide_layouts[1]  # Layout della slide (1 è per il titolo e contenuto)
        slide = presentation.slides.add_slide(slide_layout)

        # Aggiungi il titolo alla slide
        title_placeholder = slide.shapes.title
        title_placeholder.text = title

        # Aggiungi il sottotitolo (se presente)
        if subtitle:
            subtitle_placeholder = slide.placeholders[1]
            subtitle_placeholder.text = subtitle

        # Aggiungi il contenuto alla slide
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame

        for item in content:
            for key, value in item.items():
                p = text_frame.add_paragraph()
                p.text = value

        # Aggiungi le note del presentatore (se presenti)
        if presenter_notes:
            notes_slide = slide.notes_slide
            notes_text_frame = notes_slide.notes_text_frame
            notes_text_frame.text = presenter_notes

    # Salva la presentazione PowerPoint
    presentation.save('output_presentation.pptx')

    print("Presentazione creata con successo e salvata come 'output_presentation.pptx'")

    # save the output into binary form
    ppt_bytes = convert_ppt_to_bytes(presentation)

    # Usa session_state per memorizzare il buffer PPTX
    st.session_state.ppt_bytes = ppt_bytes




# Funzione per convertire il file PPTX in un buffer di byte per il download
def convert_ppt_to_bytes(prs):
    ppt_bytes = BytesIO()
    prs.save(ppt_bytes)
    ppt_bytes.seek(0)
    return ppt_bytes
