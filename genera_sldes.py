import json
from pptx import Presentation
from pptx.util import Inches, Pt

# Carica il file JSON
with open('slides.json', 'r') as f:
    data = json.load(f)


def genera_Slides():
    # Crea una nuova presentazione
    prs = Presentation()

    # Aggiungi titolo della presentazione
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = data['presentation_title']
    subtitle.text = f"{data['instructor']} - {data['date']}"

    # Aggiungi le slide definite nel file JSON
    for slide_data in data['slides']:
        # Seleziona layout della slide
        slide_layout = prs.slide_layouts[1]  # Layout con titolo e contenuto
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title

        # Aggiungi il titolo della slide
        title.text = slide_data['slide_title']

        # Aggiungi i punti elenco, se presenti
        if 'bullet_points' in slide_data:
            content = slide.shapes.placeholders[1].text_frame
            for point in slide_data['bullet_points']:
                p = content.add_paragraph()
                p.text = point

        # Aggiungi sottotitolo se presente e se esiste un segnaposto per il sottotitolo
        if 'subtitle' in slide_data and slide_data['subtitle']:
            try:
                subtitle = slide.placeholders[2]
                subtitle.text = slide_data['subtitle']
            except KeyError:
                # Se non c'è un segnaposto per il sottotitolo, lo ignoriamo
                pass

    # Salva la presentazione PowerPoint
    prs.save('Presentazione_Sistemi_Informativi_Aziendali.pptx')
